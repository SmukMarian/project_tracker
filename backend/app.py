import hashlib
import json
import logging
from datetime import date, datetime
from pathlib import Path
from typing import Optional

from io import BytesIO
import shutil

from fastapi import Depends, FastAPI, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from sqlalchemy import or_
from sqlalchemy.orm import Session, joinedload

from PIL import Image
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

from .config import get_workspace_path, set_workspace_path
from . import models, schemas
from .database import Base, SessionLocal, engine
from .logging_utils import setup_logging

Base.metadata.create_all(bind=engine)

app = FastAPI(title="Haier Project Tracker API")

LOGGER = setup_logging(get_workspace_path() / "logs")

DIST_DIR = Path(__file__).resolve().parents[1] / "frontend" / "dist"
if DIST_DIR.exists():
    app.mount("/app", StaticFiles(directory=DIST_DIR, html=True), name="spa")
    app.mount("/assets", StaticFiles(directory=DIST_DIR / "assets"), name="assets")

    @app.get("/")
    def serve_spa_root():
        return FileResponse(DIST_DIR / "index.html")

    @app.get("/app/", response_class=HTMLResponse)
    def serve_spa_index():
        return (DIST_DIR / "index.html").read_text(encoding="utf-8")

    @app.get("/cache-admin.html", response_class=HTMLResponse)
    def serve_cache_admin():
        cache_page = DIST_DIR / "cache-admin.html"
        if not cache_page.exists():
            raise HTTPException(status_code=404, detail="Cache admin page not found")
        return cache_page.read_text(encoding="utf-8")


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def _media_root() -> Path:
    workspace = get_workspace_path()
    media_root = workspace / "media"
    media_root.mkdir(parents=True, exist_ok=True)
    return media_root


def _resolve_workspace_file(path: str) -> Path:
    workspace = get_workspace_path()
    normalized = Path(path.replace("\\", "/"))
    target = (workspace / normalized).resolve()
    try:
        target.relative_to(workspace)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid file path")
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return target


@app.get("/health")
def healthcheck():
    return {"status": "ok"}


@app.get("/workspace", response_model=schemas.WorkspaceState)
def get_workspace():
    path = get_workspace_path()
    return schemas.WorkspaceState(path=str(path))


@app.post("/workspace", response_model=schemas.WorkspaceState)
def update_workspace(payload: schemas.WorkspaceUpdate):
    path = set_workspace_path(payload.path)
    setup_logging(path / "logs")
    LOGGER.info("Workspace updated to %s", path)
    return schemas.WorkspaceState(path=str(path))


@app.get("/updates/manifest", response_model=schemas.UpdateManifest)
def get_update_manifest():
    workspace = get_workspace_path()
    manifest_path = workspace / "updates" / "manifest.json"
    if not manifest_path.exists():
        raise HTTPException(status_code=404, detail="Update manifest not found")
    try:
        data = manifest_path.read_text(encoding="utf-8")
        return schemas.UpdateManifest.parse_raw(data)
    except ValueError:
        raise HTTPException(status_code=400, detail="Manifest file is not valid JSON")


@app.post("/updates/manifest", response_model=schemas.UpdateManifest)
def set_update_manifest(manifest: schemas.UpdateManifest):
    workspace = get_workspace_path()
    updates_dir = workspace / "updates"
    updates_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = updates_dir / "manifest.json"
    manifest_path.write_text(manifest.json(ensure_ascii=False, indent=2), encoding="utf-8")
    LOGGER.info("Update manifest saved to %s", manifest_path)
    return manifest


@app.post("/updates/package")
def upload_update_package(file: UploadFile = File(...)):
    workspace = get_workspace_path()
    updates_dir = workspace / "updates"
    updates_dir.mkdir(parents=True, exist_ok=True)
    dest = updates_dir / file.filename
    with dest.open("wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    digest = hashlib.sha256(dest.read_bytes()).hexdigest()
    LOGGER.info("Uploaded update package to %s (sha256=%s)", dest, digest)
    return {
        "filename": file.filename,
        "url": f"/updates/download/{file.filename}",
        "sha256": digest,
    }


@app.get("/updates/download/{filename}")
def download_update_package(filename: str):
    workspace = get_workspace_path()
    file_path = workspace / "updates" / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path)


@app.get("/files/{file_path:path}")
def serve_workspace_file(file_path: str):
    """Expose files stored under the current workspace (covers, media, attachments)."""
    target = _resolve_workspace_file(file_path)
    return FileResponse(target, headers={"Cache-Control": "no-store"})


def _cache_root() -> Path:
    workspace = get_workspace_path()
    cache_root = workspace / "cache"
    cache_root.mkdir(parents=True, exist_ok=True)
    return cache_root


def _cache_meta_path() -> Path:
    return _cache_root() / "meta.json"


def _read_cache_meta() -> dict:
    meta_path = _cache_meta_path()
    if meta_path.exists():
        try:
            return json.loads(meta_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            return {}
    return {}


def _write_cache_meta(data: dict) -> None:
    meta_path = _cache_meta_path()
    meta_path.parent.mkdir(parents=True, exist_ok=True)
    meta_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _iter_media_files():
    media_root = _media_root()
    if not media_root.exists():
        return []
    return [
        path
        for path in media_root.rglob("*")
        if path.is_file() and path.suffix.lower() in {".jpg", ".jpeg", ".png", ".webp", ".bmp"}
    ]


def _build_thumbnails(max_size: tuple[int, int] = (512, 512)) -> dict:
    media_files = _iter_media_files()
    cache_root = _cache_root()
    thumb_root = cache_root / "thumbnails"
    created = 0
    skipped = 0
    failures: list[str] = []
    for src in media_files:
        rel = src.relative_to(get_workspace_path())
        dest = thumb_root / rel
        dest_suffix = dest.suffix.lower() or ".jpg"
        format_name = "PNG" if dest_suffix == ".png" else "JPEG"
        dest_file = dest.with_suffix(".jpg" if format_name == "JPEG" else dest.suffix)
        if dest_file.exists() and dest_file.stat().st_mtime >= src.stat().st_mtime:
            skipped += 1
            continue
        dest_file.parent.mkdir(parents=True, exist_ok=True)
        try:
            with Image.open(src) as img:
                img = img.convert("RGB")
                img.thumbnail(max_size)
                img.save(dest_file, format=format_name, optimize=True, quality=85)
                created += 1
        except Exception as exc:  # pragma: no cover - defensive log
            failures.append(f"{src.name}: {exc}")
    meta = _read_cache_meta()
    meta["last_thumbnail_run"] = datetime.utcnow().isoformat() + "Z"
    _write_cache_meta(meta)
    return {"processed": len(media_files), "created": created, "skipped": skipped, "failures": failures}


def _prewarm_cache() -> dict:
    files = list(_media_root().rglob("*"))
    touched = 0
    for path in files:
        if not path.is_file():
            continue
        try:
            with path.open("rb") as handle:
                handle.read(1024)
            touched += 1
        except OSError:
            continue
    meta = _read_cache_meta()
    meta["last_prewarm"] = datetime.utcnow().isoformat() + "Z"
    _write_cache_meta(meta)
    return {"files_touched": touched, "total_items": len(files)}


def _cache_status() -> dict:
    media_files = _iter_media_files()
    thumb_root = _cache_root() / "thumbnails"
    thumbnails = list(thumb_root.rglob("*")) if thumb_root.exists() else []
    meta = _read_cache_meta()
    return {
        "media_files": len(media_files),
        "thumbnails": len([p for p in thumbnails if p.is_file()]),
        "last_thumbnail_run": meta.get("last_thumbnail_run"),
        "last_prewarm": meta.get("last_prewarm"),
    }


@app.get("/cache/status", response_model=schemas.CacheStatus)
def cache_status():
    return _cache_status()


@app.post("/cache/thumbnails", response_model=schemas.CacheRunResult)
def cache_build_thumbnails():
    return _build_thumbnails()


@app.post("/cache/prewarm", response_model=schemas.CachePrewarmResult)
def cache_prewarm():
    return _prewarm_cache()


@app.delete("/cache", status_code=204)
def cache_clear():
    cache_root = _cache_root()
    try:
        shutil.rmtree(cache_root)
    except FileNotFoundError:
        return
    cache_root.mkdir(parents=True, exist_ok=True)


def _status_value(status: str, inprogress_coeff: float) -> float:
    mapping = {
        "done": 1.0,
        "in_progress": inprogress_coeff,
        "blocked": 0.0,
        "todo": 0.0,
    }
    return mapping.get(status, 0.0)


def _compute_step_progress(step: models.Step, inprogress_coeff: float) -> int:
    if step.subtasks:
        total_weight = sum(subtask.weight for subtask in step.subtasks)
        if total_weight == 0:
            return 0
        value = sum(
            _status_value(subtask.status, inprogress_coeff) * subtask.weight
            for subtask in step.subtasks
        )
        return round((value / total_weight) * 100)
    return round(_status_value(step.status, inprogress_coeff) * 100)


def _apply_progress(project: models.Project) -> models.Project:
    step_progresses: list[tuple[int, float]] = []
    steps_total = 0
    steps_done = 0
    subtasks_total = 0
    subtasks_done = 0
    for step in project.steps:
        step.progress_percent = _compute_step_progress(step, project.inprogress_coeff)
        weight = step.weight or 0.0
        step_progresses.append((step.progress_percent, weight if weight > 0 else 0))
        steps_total += 1
        if step.progress_percent >= 100:
            steps_done += 1
        for subtask in step.subtasks:
            subtasks_total += 1
            if subtask.status == schemas.TaskStatus.DONE:
                subtasks_done += 1
    total_weight = sum(weight for _, weight in step_progresses)
    if total_weight > 0:
        weighted_progress = sum(progress * weight for progress, weight in step_progresses)
        project.progress_percent = round(weighted_progress / total_weight)
    elif step_progresses:
        project.progress_percent = round(
            sum(progress for progress, _ in step_progresses) / len(step_progresses)
        )
    else:
        project.progress_percent = 0
    project.steps_total = steps_total
    project.steps_done = steps_done
    project.subtasks_total = subtasks_total
    project.subtasks_done = subtasks_done
    return project


def _apply_category_metrics(category: models.Category, db: Session) -> models.Category:
    projects = [
        _apply_progress(project)
        for project in category.projects
    ]
    if projects:
        category.progress_percent = round(
            sum(project.progress_percent for project in projects) / len(projects)
        )
    else:
        category.progress_percent = 0
    category.average_progress = category.progress_percent
    category.kpi = _kpi_report(db, category.id)
    return category


def _workbook_response(wb: Workbook, filename: str) -> StreamingResponse:
    buffer = BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


def _binary_response(buffer: BytesIO, media_type: str, filename: str) -> StreamingResponse:
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


def _export_categories_excel(db: Session) -> StreamingResponse:
    categories = (
        db.query(models.Category)
        .options(
            joinedload(models.Category.projects)
            .joinedload(models.Project.steps)
            .joinedload(models.Step.subtasks)
        )
        .order_by(models.Category.name)
        .all()
    )

    wb = Workbook()
    ws_categories = wb.active
    ws_categories.title = "Categories"
    ws_categories.append(["ID", "Название", "Проектов"])
    for category in categories:
        ws_categories.append([category.id, category.name, len(category.projects)])

    ws_projects = wb.create_sheet("Projects")
    ws_projects.append(
        [
            "ID",
            "Название",
            "Код",
            "Статус",
            "Категория",
            "PM",
            "Дата старта",
            "Целевая дата",
            "Прогресс %",
        ]
    )
    for category in categories:
        for project in category.projects:
            _apply_progress(project)
            ws_projects.append(
                [
                    project.id,
                    project.name,
                    project.code,
                    project.status,
                    category.name,
                    project.owner_id,
                    project.start_date,
                    project.target_date,
                    project.progress_percent,
                ]
            )

    return _workbook_response(wb, "categories.xlsx")


def _kpi_report(db: Session, category_id: Optional[int] = None) -> schemas.KPIReport:
    projects_query = db.query(models.Project).options(
        joinedload(models.Project.steps).joinedload(models.Step.subtasks)
    )
    if category_id is not None:
        projects_query = projects_query.filter(models.Project.category_id == category_id)
    projects = projects_query.all()

    total_projects = len(projects)
    active_projects = len([p for p in projects if p.status == schemas.ProjectStatus.ACTIVE])
    archived_projects = len([p for p in projects if p.status == schemas.ProjectStatus.ARCHIVED])

    total_progress = 0
    steps_total = 0
    steps_done = 0
    subtasks_total = 0
    subtasks_done = 0

    for project in projects:
        _apply_progress(project)
        total_progress += project.progress_percent
        steps_total += project.steps_total
        steps_done += project.steps_done
        subtasks_total += project.subtasks_total
        subtasks_done += project.subtasks_done

    average_progress = round(total_progress / total_projects, 2) if total_projects else 0.0
    return schemas.KPIReport(
        total_projects=total_projects,
        active_projects=active_projects,
        archived_projects=archived_projects,
        average_progress=average_progress,
        steps_total=steps_total,
        steps_done=steps_done,
        subtasks_total=subtasks_total,
        subtasks_done=subtasks_done,
    )


def _export_projects_word(db: Session, category_id: Optional[int] = None) -> StreamingResponse:
    query = (
        db.query(models.Project)
        .options(
            joinedload(models.Project.category),
            joinedload(models.Project.steps).joinedload(models.Step.subtasks),
        )
        .order_by(models.Project.start_date.is_(None), models.Project.start_date)
    )
    if category_id is not None:
        query = query.filter(models.Project.category_id == category_id)
    projects = query.all()

    doc = Document()
    doc.add_heading("Проекты Haier", level=0)
    if category_id is not None:
        category = db.query(models.Category).filter(models.Category.id == category_id).first()
        if category:
            doc.add_paragraph(f"Категория: {category.name}")

    for project in projects:
        _apply_progress(project)
        doc.add_heading(f"{project.name} ({project.code or 'без кода'})", level=1)
        doc.add_paragraph(f"Статус: {project.status}")
        if project.owner_id:
            doc.add_paragraph(f"PM: {project.owner_id}")
        if project.start_date or project.target_date:
            doc.add_paragraph(f"Период: {project.start_date or '-'} → {project.target_date or '-'}")
        doc.add_paragraph(f"Прогресс: {project.progress_percent}%")
        if project.description:
            doc.add_paragraph(project.description)
        if project.steps:
            doc.add_paragraph("Шаги:")
            for step in project.steps:
                doc.add_paragraph(
                    f"- {step.name} — {step.status} ({step.progress_percent}% по шагу)", style="List Bullet"
                )

    stream = BytesIO()
    doc.save(stream)
    return _binary_response(stream, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "projects.docx")


def _export_category_presentation(db: Session, category_id: int) -> StreamingResponse:
    category = (
        db.query(models.Category)
        .options(joinedload(models.Category.projects).joinedload(models.Project.steps))
        .filter(models.Category.id == category_id)
        .first()
    )
    if not category:
        raise HTTPException(status_code=404, detail="Category not found")

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"Категория: {category.name}"
    subtitle = slide.placeholders[1]
    subtitle.text = "Презентация проектов"

    bullet_layout = prs.slide_layouts[1]
    for project in category.projects:
        _apply_progress(project)
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = project.name
        body = slide.shapes.placeholders[1].text_frame
        body.text = f"Код: {project.code or '-'}"
        body.add_paragraph().text = f"Статус: {project.status}"
        body.add_paragraph().text = f"PM: {project.owner_id or '-'}"
        body.add_paragraph().text = f"Прогресс: {project.progress_percent}%"
        if project.target_date:
            body.add_paragraph().text = f"Целевая дата: {project.target_date}"
        if project.steps:
            steps_paragraph = body.add_paragraph("Шаги:")
            steps_paragraph.level = 0
            for step in project.steps:
                p = body.add_paragraph(f"• {step.name} — {step.status}")
                p.level = 1

    stream = BytesIO()
    prs.save(stream)
    return _binary_response(stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "category.pptx")


def _load_project_with_steps(project_id: int, db: Session) -> models.Project:
    project = (
        db.query(models.Project)
        .options(joinedload(models.Project.steps).joinedload(models.Step.subtasks))
        .filter(models.Project.id == project_id)
        .first()
    )
    if project:
        _apply_progress(project)
    return project


@app.post("/categories", response_model=schemas.Category)
def create_category(category: schemas.CategoryCreate, db: Session = Depends(get_db)):
    existing = db.query(models.Category).filter(models.Category.name == category.name).first()
    if existing:
        raise HTTPException(status_code=400, detail="Category already exists")
    db_category = models.Category(name=category.name)
    db.add(db_category)
    db.commit()
    db.refresh(db_category)
    return db_category


@app.get("/categories", response_model=list[schemas.Category])
def list_categories(db: Session = Depends(get_db)):
    categories = (
        db.query(models.Category)
        .options(
            joinedload(models.Category.projects)
            .joinedload(models.Project.steps)
            .joinedload(models.Step.subtasks)
        )
        .order_by(models.Category.name)
        .all()
    )
    return [_apply_category_metrics(category, db) for category in categories]


@app.get("/export/categories/excel")
def export_categories_excel(db: Session = Depends(get_db)):
    return _export_categories_excel(db)


@app.get("/export/projects/word")
def export_projects_word(category_id: Optional[int] = Query(None), db: Session = Depends(get_db)):
    return _export_projects_word(db, category_id)


@app.get("/export/category/{category_id}/presentation")
def export_category_presentation(category_id: int, db: Session = Depends(get_db)):
    return _export_category_presentation(db, category_id)


@app.get("/kpi", response_model=schemas.KPIReport)
def kpi_report(category_id: Optional[int] = Query(None), db: Session = Depends(get_db)):
    return _kpi_report(db, category_id)


@app.post("/pms", response_model=schemas.PM)
def create_pm(pm: schemas.PMCreate, db: Session = Depends(get_db)):
    existing = db.query(models.PM).filter(models.PM.name == pm.name).first()
    if existing:
        raise HTTPException(status_code=400, detail="PM already exists")
    db_pm = models.PM(name=pm.name)
    db.add(db_pm)
    db.commit()
    db.refresh(db_pm)
    return db_pm


@app.get("/pms", response_model=list[schemas.PM])
def list_pms(db: Session = Depends(get_db)):
    return db.query(models.PM).all()


@app.post("/projects", response_model=schemas.Project)
def create_project(project: schemas.ProjectCreate, db: Session = Depends(get_db)):
    db_project = models.Project(**project.dict())
    db.add(db_project)
    db.commit()
    db.refresh(db_project)
    db.refresh(db_project, attribute_names=["steps", "characteristics", "attachments"])
    return _apply_progress(db_project)


@app.post("/projects/status", response_model=list[schemas.Project])
def bulk_update_project_status(
    payload: schemas.BulkProjectStatusUpdate, db: Session = Depends(get_db)
):
    projects = (
        db.query(models.Project)
        .options(
            joinedload(models.Project.steps).joinedload(models.Step.subtasks),
            joinedload(models.Project.characteristics),
            joinedload(models.Project.attachments),
        )
        .filter(models.Project.id.in_(payload.ids))
        .all()
    )
    if not projects:
        raise HTTPException(status_code=404, detail="No matching projects found")
    for project in projects:
        project.status = payload.status.value
    db.commit()
    for project in projects:
        db.refresh(project)
    return [_apply_progress(project) for project in projects]


@app.patch("/projects/{project_id}", response_model=schemas.Project)
def update_project(project_id: int, update: schemas.ProjectUpdate, db: Session = Depends(get_db)):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    for key, value in update.dict(exclude_unset=True).items():
        setattr(project, key, value)
    db.commit()
    db.refresh(project)
    db.refresh(project, attribute_names=["steps", "characteristics", "attachments"])
    return _apply_progress(project)


@app.get("/projects/{project_id}", response_model=schemas.Project)
def get_project(project_id: int, db: Session = Depends(get_db)):
    project = (
        db.query(models.Project)
        .options(
            joinedload(models.Project.steps).joinedload(models.Step.subtasks),
            joinedload(models.Project.characteristics),
            joinedload(models.Project.attachments),
        )
        .filter(models.Project.id == project_id)
        .first()
    )
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    return _apply_progress(project)


@app.get("/projects", response_model=list[schemas.Project])
def list_projects(
    category_id: Optional[int] = None,
    owner_id: Optional[int] = None,
    status: Optional[schemas.ProjectStatus] = None,
    search: Optional[str] = Query(None, description="Search by name, code, or status"),
    db: Session = Depends(get_db),
):
    query = (
        db.query(models.Project)
        .options(
            joinedload(models.Project.steps).joinedload(models.Step.subtasks),
            joinedload(models.Project.characteristics),
            joinedload(models.Project.attachments),
        )
        .order_by(models.Project.id)
    )

    if category_id is not None:
        query = query.filter(models.Project.category_id == category_id)
    if owner_id is not None:
        query = query.filter(models.Project.owner_id == owner_id)
    if status is not None:
        query = query.filter(models.Project.status == status.value)
    if search:
        like = f"%{search}%"
        query = query.filter(
            or_(
                models.Project.name.ilike(like),
                models.Project.code.ilike(like),
                models.Project.status.ilike(like),
            )
        )

    projects = query.all()
    return [_apply_progress(project) for project in projects]


@app.post("/projects/bulk-delete", status_code=204)
def bulk_delete_projects(payload: schemas.BulkDeleteRequest, db: Session = Depends(get_db)):
    projects = db.query(models.Project).filter(models.Project.id.in_(payload.ids)).all()
    if not projects:
        raise HTTPException(status_code=404, detail="No matching projects found")
    for project in projects:
        db.delete(project)
    db.commit()


@app.get("/projects/{project_id}/steps", response_model=list[schemas.Step])
def list_steps(
    project_id: int,
    status: Optional[schemas.TaskStatus] = None,
    assignee_id: Optional[int] = None,
    search: Optional[str] = Query(None, description="Search by name or description"),
    db: Session = Depends(get_db),
):
    steps = (
        db.query(models.Step)
        .options(joinedload(models.Step.subtasks))
        .filter(models.Step.project_id == project_id)
        .order_by(models.Step.order_index, models.Step.id)
    )
    if status:
        steps = steps.filter(models.Step.status == status.value)
    if assignee_id:
        steps = steps.filter(models.Step.assignee_id == assignee_id)
    if search:
        like = f"%{search}%"
        steps = steps.filter(or_(models.Step.name.ilike(like), models.Step.description.ilike(like)))

    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    step_rows = steps.all()
    for step in step_rows:
        step.progress_percent = _compute_step_progress(step, project.inprogress_coeff)
    return step_rows


@app.post("/steps", response_model=schemas.Step)
def create_step(step: schemas.StepCreate, db: Session = Depends(get_db)):
    project = db.query(models.Project).filter(models.Project.id == step.project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    db_step = models.Step(**step.dict())
    db.add(db_step)
    db.commit()
    db.refresh(db_step)
    db.refresh(db_step, attribute_names=["subtasks"])
    db_step.progress_percent = _compute_step_progress(db_step, project.inprogress_coeff)
    return db_step


@app.patch("/steps/{step_id}", response_model=schemas.Step)
def update_step(step_id: int, update: schemas.StepUpdate, db: Session = Depends(get_db)):
    step = db.query(models.Step).filter(models.Step.id == step_id).first()
    if step is None:
        raise HTTPException(status_code=404, detail="Step not found")
    for key, value in update.dict(exclude_unset=True).items():
        setattr(step, key, value)
    db.commit()
    db.refresh(step)
    db.refresh(step, attribute_names=["subtasks"])
    project = db.query(models.Project).filter(models.Project.id == step.project_id).first()
    step.progress_percent = _compute_step_progress(step, project.inprogress_coeff)
    return step


@app.delete("/steps/{step_id}", status_code=204)
def delete_step(step_id: int, db: Session = Depends(get_db)):
    step = db.query(models.Step).filter(models.Step.id == step_id).first()
    if step is None:
        raise HTTPException(status_code=404, detail="Step not found")
    db.delete(step)
    db.commit()


@app.post("/steps/bulk-delete", status_code=204)
def bulk_delete_steps(payload: schemas.BulkDeleteRequest, db: Session = Depends(get_db)):
    steps = db.query(models.Step).filter(models.Step.id.in_(payload.ids)).all()
    if not steps:
        raise HTTPException(status_code=404, detail="No matching steps found")
    for step in steps:
        db.delete(step)
    db.commit()


@app.post("/projects/{project_id}/steps/reorder", response_model=list[schemas.Step])
def reorder_steps(project_id: int, payload: schemas.OrderUpdate, db: Session = Depends(get_db)):
    step_ids = payload.ids
    steps = db.query(models.Step).filter(models.Step.project_id == project_id).all()
    existing_ids = {step.id for step in steps}
    if set(step_ids) - existing_ids:
        raise HTTPException(status_code=400, detail="One or more steps do not belong to the project")

    order_map = {step_id: index for index, step_id in enumerate(step_ids)}
    for step in steps:
        if step.id in order_map:
            step.order_index = order_map[step.id]
    db.commit()
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    for step in steps:
        step.progress_percent = _compute_step_progress(step, project.inprogress_coeff)
    return sorted(steps, key=lambda s: (s.order_index, s.id))


@app.post("/subtasks", response_model=schemas.Subtask)
def create_subtask(subtask: schemas.SubtaskCreate, db: Session = Depends(get_db)):
    step = db.query(models.Step).filter(models.Step.id == subtask.step_id).first()
    if step is None:
        raise HTTPException(status_code=404, detail="Step not found")
    db_subtask = models.Subtask(**subtask.dict())
    db.add(db_subtask)
    db.commit()
    db.refresh(db_subtask)
    return db_subtask


@app.patch("/subtasks/{subtask_id}", response_model=schemas.Subtask)
def update_subtask(subtask_id: int, update: schemas.SubtaskUpdate, db: Session = Depends(get_db)):
    subtask = db.query(models.Subtask).filter(models.Subtask.id == subtask_id).first()
    if subtask is None:
        raise HTTPException(status_code=404, detail="Subtask not found")
    for key, value in update.dict(exclude_unset=True).items():
        setattr(subtask, key, value)
    db.commit()
    db.refresh(subtask)
    return subtask


@app.delete("/subtasks/{subtask_id}", status_code=204)
def delete_subtask(subtask_id: int, db: Session = Depends(get_db)):
    subtask = db.query(models.Subtask).filter(models.Subtask.id == subtask_id).first()
    if subtask is None:
        raise HTTPException(status_code=404, detail="Subtask not found")
    db.delete(subtask)
    db.commit()


@app.post("/subtasks/bulk-delete", status_code=204)
def bulk_delete_subtasks(payload: schemas.BulkDeleteRequest, db: Session = Depends(get_db)):
    subtasks = db.query(models.Subtask).filter(models.Subtask.id.in_(payload.ids)).all()
    if not subtasks:
        raise HTTPException(status_code=404, detail="No matching subtasks found")
    for subtask in subtasks:
        db.delete(subtask)
    db.commit()


@app.post("/steps/{step_id}/subtasks/reorder", response_model=list[schemas.Subtask])
def reorder_subtasks(step_id: int, payload: schemas.OrderUpdate, db: Session = Depends(get_db)):
    subtask_ids = payload.ids
    subtasks = db.query(models.Subtask).filter(models.Subtask.step_id == step_id).all()
    existing_ids = {subtask.id for subtask in subtasks}
    if set(subtask_ids) - existing_ids:
        raise HTTPException(status_code=400, detail="One or more subtasks do not belong to the step")

    order_map = {subtask_id: index for index, subtask_id in enumerate(subtask_ids)}
    for subtask in subtasks:
        if subtask.id in order_map:
            subtask.order_index = order_map[subtask.id]
    db.commit()
    return sorted(subtasks, key=lambda s: (s.order_index, s.id))


@app.get("/steps/{step_id}/subtasks", response_model=list[schemas.Subtask])
def list_subtasks(
    step_id: int,
    status: Optional[schemas.TaskStatus] = None,
    search: Optional[str] = Query(None, description="Search by name"),
    db: Session = Depends(get_db),
):
    step = db.query(models.Step).filter(models.Step.id == step_id).first()
    if step is None:
        raise HTTPException(status_code=404, detail="Step not found")
    subtasks = db.query(models.Subtask).filter(models.Subtask.step_id == step_id)
    if status:
        subtasks = subtasks.filter(models.Subtask.status == status.value)
    if search:
        like = f"%{search}%"
        subtasks = subtasks.filter(models.Subtask.name.ilike(like))
    return subtasks.order_by(models.Subtask.order_index, models.Subtask.id).all()


@app.get("/projects/{project_id}/characteristics", response_model=list[schemas.ProjectCharacteristic])
def list_characteristics(project_id: int, db: Session = Depends(get_db)):
    return db.query(models.ProjectCharacteristic).filter(models.ProjectCharacteristic.project_id == project_id).all()


def _replace_characteristics(project_id: int, items: list[schemas.ProjectCharacteristicBase], db: Session):
    db.query(models.ProjectCharacteristic).filter(models.ProjectCharacteristic.project_id == project_id).delete()
    for item in items:
        db.add(
            models.ProjectCharacteristic(
                project_id=project_id,
                parameter=item.parameter,
                value=item.value,
            )
        )
    db.commit()
    return (
        db.query(models.ProjectCharacteristic)
        .filter(models.ProjectCharacteristic.project_id == project_id)
        .order_by(models.ProjectCharacteristic.id)
        .all()
    )


@app.post(
    "/projects/{project_id}/characteristics/import/json", response_model=list[schemas.ProjectCharacteristic]
)
def import_characteristics_json(
    project_id: int, payload: schemas.ProjectCharacteristicImport, db: Session = Depends(get_db)
):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    return _replace_characteristics(project_id, payload.items, db)


@app.get(
    "/projects/{project_id}/characteristics/export/json", response_model=list[schemas.ProjectCharacteristic]
)
def export_characteristics_json(project_id: int, db: Session = Depends(get_db)):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    return (
        db.query(models.ProjectCharacteristic)
        .filter(models.ProjectCharacteristic.project_id == project_id)
        .order_by(models.ProjectCharacteristic.id)
        .all()
    )


@app.post(
    "/projects/{project_id}/characteristics/import/excel", response_model=list[schemas.ProjectCharacteristic]
)
async def import_characteristics_excel(
    project_id: int, file: UploadFile = File(...), db: Session = Depends(get_db)
):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    content = await file.read()
    workbook = load_workbook(filename=BytesIO(content))
    sheet = workbook.active
    items: list[schemas.ProjectCharacteristicBase] = []
    for row in sheet.iter_rows(min_row=2, values_only=True):
        if not row or not row[0]:
            continue
        parameter = str(row[0])
        value = "" if len(row) < 2 or row[1] is None else str(row[1])
        items.append(schemas.ProjectCharacteristicBase(parameter=parameter, value=value))
    return _replace_characteristics(project_id, items, db)


@app.get("/projects/{project_id}/characteristics/export/excel")
def export_characteristics_excel(project_id: int, db: Session = Depends(get_db)):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")

    rows = (
        db.query(models.ProjectCharacteristic)
        .filter(models.ProjectCharacteristic.project_id == project_id)
        .order_by(models.ProjectCharacteristic.id)
        .all()
    )
    wb = Workbook()
    sheet = wb.active
    sheet.title = "Characteristics"
    sheet.append(["Параметр", "Значение"])
    for row in rows:
        sheet.append([row.parameter, row.value])
    return _workbook_response(wb, f"project_{project_id}_characteristics.xlsx")


@app.post("/characteristics", response_model=schemas.ProjectCharacteristic)
def create_characteristic(
    characteristic: schemas.ProjectCharacteristicCreate, db: Session = Depends(get_db)
):
    project = db.query(models.Project).filter(models.Project.id == characteristic.project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")
    db_characteristic = models.ProjectCharacteristic(**characteristic.dict())
    db.add(db_characteristic)
    db.commit()
    db.refresh(db_characteristic)
    return db_characteristic


@app.patch("/characteristics/{characteristic_id}", response_model=schemas.ProjectCharacteristic)
def update_characteristic(
    characteristic_id: int, update: schemas.ProjectCharacteristicUpdate, db: Session = Depends(get_db)
):
    characteristic = (
        db.query(models.ProjectCharacteristic)
        .filter(models.ProjectCharacteristic.id == characteristic_id)
        .first()
    )
    if characteristic is None:
        raise HTTPException(status_code=404, detail="Characteristic not found")
    for key, value in update.dict(exclude_unset=True).items():
        setattr(characteristic, key, value)
    db.commit()
    db.refresh(characteristic)
    return characteristic


@app.delete("/characteristics/{characteristic_id}", status_code=204)
def delete_characteristic(characteristic_id: int, db: Session = Depends(get_db)):
    characteristic = (
        db.query(models.ProjectCharacteristic)
        .filter(models.ProjectCharacteristic.id == characteristic_id)
        .first()
    )
    if characteristic is None:
        raise HTTPException(status_code=404, detail="Characteristic not found")
    db.delete(characteristic)
    db.commit()


@app.post("/attachments", response_model=schemas.Attachment)
def create_attachment(attachment: schemas.AttachmentCreate, db: Session = Depends(get_db)):
    if attachment.project_id is None and attachment.step_id is None:
        raise HTTPException(status_code=400, detail="Attachment must reference a project or step")
    db_attachment = models.Attachment(**attachment.dict())
    db.add(db_attachment)
    db.commit()
    db.refresh(db_attachment)
    LOGGER.info(
        "Attachment record created: path=%s project_id=%s step_id=%s",
        db_attachment.path,
        db_attachment.project_id,
        db_attachment.step_id,
    )
    return db_attachment


def _ensure_target_exists(db: Session, project_id: Optional[int], step_id: Optional[int]) -> None:
    if project_id is None and step_id is None:
        raise HTTPException(status_code=400, detail="Attachment must reference a project or step")

    if project_id is not None:
        exists = db.query(models.Project.id).filter(models.Project.id == project_id).first()
        if not exists:
            raise HTTPException(status_code=404, detail="Project not found")
    if step_id is not None:
        exists = db.query(models.Step.id).filter(models.Step.id == step_id).first()
        if not exists:
            raise HTTPException(status_code=404, detail="Step not found")


def _persist_file(file: UploadFile, destination: Path) -> Path:
    original_name = Path(file.filename or "attachment")
    safe_name = original_name.name
    target = destination / safe_name
    counter = 1
    while target.exists():
        target = destination / f"{original_name.stem}_{counter}{original_name.suffix}"
        counter += 1

    with target.open("wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    return target


@app.post("/attachments/upload", response_model=schemas.Attachment)
async def upload_attachment(
    project_id: Optional[int] = Form(None),
    step_id: Optional[int] = Form(None),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    _ensure_target_exists(db, project_id, step_id)

    media_root = _media_root()
    target_dir = media_root / f"project_{project_id}" if project_id else media_root / f"step_{step_id}"
    target_dir.mkdir(parents=True, exist_ok=True)

    stored_file = _persist_file(file, target_dir)
    relative_path = stored_file.relative_to(get_workspace_path())

    LOGGER.info(
        "Attachment uploaded: filename=%s stored_as=%s project_id=%s step_id=%s",
        file.filename,
        relative_path,
        project_id,
        step_id,
    )

    db_attachment = models.Attachment(
        path=str(relative_path),
        project_id=project_id,
        step_id=step_id,
        added_at=date.today(),
    )
    db.add(db_attachment)
    db.commit()
    db.refresh(db_attachment)
    return db_attachment


@app.delete("/attachments/{attachment_id}", status_code=204)
def delete_attachment(attachment_id: int, db: Session = Depends(get_db)):
    attachment = db.query(models.Attachment).filter(models.Attachment.id == attachment_id).first()
    if attachment is None:
        raise HTTPException(status_code=404, detail="Attachment not found")
    workspace = get_workspace_path()
    file_path = workspace / attachment.path
    if file_path.exists():
        try:
            file_path.unlink()
        except OSError:
            LOGGER.warning("Failed to delete attachment file %s", file_path)
    LOGGER.info(
        "Attachment removed: id=%s path=%s project_id=%s step_id=%s",
        attachment_id,
        attachment.path,
        attachment.project_id,
        attachment.step_id,
    )
    db.delete(attachment)
    db.commit()


@app.post("/projects/{project_id}/cover", response_model=schemas.Project)
async def upload_project_cover(
    project_id: int, file: UploadFile = File(...), db: Session = Depends(get_db)
):
    project = db.query(models.Project).filter(models.Project.id == project_id).first()
    if project is None:
        raise HTTPException(status_code=404, detail="Project not found")

    media_root = _media_root()
    target_dir = media_root / f"project_{project_id}" / "cover"
    target_dir.mkdir(parents=True, exist_ok=True)

    stored_file = _persist_file(file, target_dir)
    new_relative = stored_file.relative_to(get_workspace_path())

    if project.cover_image:
        previous = get_workspace_path() / project.cover_image
        if previous.exists():
            try:
                previous.unlink()
            except OSError:
                LOGGER.warning("Failed to remove old cover %s", previous)

    project.cover_image = str(new_relative)
    db.commit()
    db.refresh(project)
    db.refresh(project, attribute_names=["steps", "characteristics", "attachments"])
    return _apply_progress(project)


@app.get("/projects/{project_id}/attachments", response_model=list[schemas.Attachment])
def list_project_attachments(project_id: int, db: Session = Depends(get_db)):
    return db.query(models.Attachment).filter(models.Attachment.project_id == project_id).all()


@app.get("/steps/{step_id}/attachments", response_model=list[schemas.Attachment])
def list_step_attachments(step_id: int, db: Session = Depends(get_db)):
    return db.query(models.Attachment).filter(models.Attachment.step_id == step_id).all()
